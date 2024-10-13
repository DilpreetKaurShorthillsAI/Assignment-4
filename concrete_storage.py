import os
import csv
import sqlite3
import fitz  # PyMuPDF for PDF handling
import docx
from pptx import Presentation

class FileLoader:
    """Abstract class for file loaders."""
    def extract_text(self):
        raise NotImplementedError

    def extract_links(self):
        raise NotImplementedError

    def extract_images(self):
        raise NotImplementedError

    def extract_tables(self):
        raise NotImplementedError


# PDF Loader
class PDFLoader(FileLoader):
    def __init__(self, filepath):
        self.filepath = filepath
        self.doc = fitz.open(filepath)

    def extract_text(self):
        text_data = []
        for page in self.doc:
            text_data.append(page.get_text())
        return text_data

    def extract_links(self):
        links_data = []
        for page_num, page in enumerate(self.doc, start=1):
            links = page.get_links()
            for link in links:
                links_data.append({
                    'text': link.get('text', ''),
                    'url': link.get('uri', ''),
                    'page': page_num
                })
        return links_data

    def extract_images(self):
        images_data = []
        for page_num, page in enumerate(self.doc, start=1):
            images = page.get_images(full=True)
            for img in images:
                xref = img[0]
                base_image = self.doc.extract_image(xref)
                image_bytes = base_image['image']
                image_format = base_image['ext']
                images_data.append({
                    'image_bytes': image_bytes,
                    'image_format': image_format,
                    'page': page_num
                })
        return images_data

    def extract_tables(self):
        # PyMuPDF does not have native table extraction support. Use OCR or manual table detection.
        # For demonstration, returning an empty list.
        return []


# DOCX Loader
class DOCXLoader(FileLoader):
    def __init__(self, filepath):
        self.doc = docx.Document(filepath)

    def extract_text(self):
        return [para.text for para in self.doc.paragraphs]

    def extract_links(self):
        links_data = []
        for rel in self.doc.part.rels:
            if "hyperlink" in rel:
                links_data.append({
                    'text': rel.target_part.text,
                    'url': rel.target_ref
                })
        return links_data

    def extract_images(self):
        images_data = []
        for rel in self.doc.part.rels.values():
            if "image" in rel.target_ref:
                image_bytes = rel.target_part.blob
                image_format = rel.target_ref.split(".")[-1]
                images_data.append({
                    'image_bytes': image_bytes,
                    'image_format': image_format
                })
        return images_data

    def extract_tables(self):
        tables = []
        for table in self.doc.tables:
            table_data = [[cell.text for cell in row.cells] for row in table.rows]
            tables.append(table_data)
        return tables


# PPT Loader
class PPTLoader(FileLoader):
    def __init__(self, filepath):
        self.prs = Presentation(filepath)

    def extract_text(self):
        text_data = []
        for slide in self.prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text_data.append(shape.text)
        return text_data

    def extract_links(self):
        links_data = []
        for slide_num, slide in enumerate(self.prs.slides, start=1):
            for shape in slide.shapes:
                if hasattr(shape, "hyperlink") and shape.hyperlink:
                    links_data.append({
                        'text': shape.text,
                        'url': shape.hyperlink.address,
                        'slide': slide_num
                    })
        return links_data

    def extract_images(self):
        images_data = []
        for slide_num, slide in enumerate(self.prs.slides, start=1):
            for shape in slide.shapes:
                if shape.shape_type == 13:  # Picture
                    image_bytes = shape.image.blob
                    image_format = shape.image.ext
                    images_data.append({
                        'image_bytes': image_bytes,
                        'image_format': image_format,
                        'slide': slide_num
                    })
        return images_data

    def extract_tables(self):
        tables = []
        for slide in self.prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "table"):
                    table_data = [[cell.text for cell in row.cells] for row in shape.table.rows]
                    tables.append(table_data)
        return tables


# File and SQL Storage Classes
class FileStorage:
    def __init__(self, directory):
        """Initialize storage directory."""
        self.directory = directory
        if not os.path.exists(self.directory):
            os.makedirs(self.directory)

    def save_text_data(self, text_data, filename):
        """Save extracted text data to CSV."""
        file_path = os.path.join(self.directory, filename)
        try:
            with open(file_path, mode='w', newline='', encoding='utf-8') as file:
                writer = csv.writer(file)

                # Handle lists of strings (e.g., text)
                if isinstance(text_data, list):
                    for row in text_data:
                        writer.writerow([row])
            print(f"Text data saved successfully to {filename}")
        except Exception as e:
            print(f"Error saving text data to {filename}: {e}")

    def save_links_data(self, links_data, filename):
        """Save extracted links data to CSV."""
        file_path = os.path.join(self.directory, filename)
        try:
            with open(file_path, mode='w', newline='', encoding='utf-8') as file:
                writer = csv.writer(file)

                # Handle dictionaries (e.g., links)
                if isinstance(links_data, list) and len(links_data) > 0 and isinstance(links_data[0], dict):
                    # Write header row based on the keys of the first dict
                    writer.writerow(links_data[0].keys())
                    for item in links_data:
                        writer.writerow(item.values())
            print(f"Links data saved successfully to {filename}")
        except Exception as e:
            print(f"Error saving links data to {filename}: {e}")

    def save_images(self, images_data):
        """Save extracted images to disk."""
        images_dir = os.path.join(self.directory, 'images')
        os.makedirs(images_dir, exist_ok=True)

        for i, img_data in enumerate(images_data):
            image_format = img_data.get('image_format', 'png')  # Default to 'png'
            file_path = os.path.join(images_dir, f'image_{i + 1}.{image_format}')
            try:
                with open(file_path, "wb") as img_file:
                    img_file.write(img_data['image_bytes'])
                print(f"Image saved successfully: {file_path}")
            except Exception as e:
                print(f"Error saving image {file_path}: {e}")

    def save_tables(self, tables_data):
        """Save extracted tables to separate CSV files."""
        tables_dir = os.path.join(self.directory, 'tables')
        os.makedirs(tables_dir, exist_ok=True)

        for i, table_data in enumerate(tables_data):
            file_path = os.path.join(tables_dir, f'table_{i + 1}.csv')

            # Check if table_data is a list (directly write rows if it's a list)
            if isinstance(table_data, list):
                # Ensure each cell is converted to a string, but preserve original data types (str, int, float, etc.)
                modified_table = [
                    [str(cell) if cell is not None else '' for cell in row] for row in table_data
                ]
            else:
                # If the structure is different, raise an error or handle accordingly
                print(f"Unexpected table format for table {i + 1}: {table_data}")
                continue

            try:
                with open(file_path, mode='w', newline='', encoding='utf-8') as file:
                    writer = csv.writer(file)
                    writer.writerows(modified_table)
                print(f"Table {i + 1} saved successfully.")
            except Exception as e:
                print(f"Error saving table {file_path}: {e}")


class SQLStorage:
    def __init__(self, db_name):
        """Initialize SQL storage with a database."""
        self.db_name = db_name
        self.connection = sqlite3.connect(self.db_name)
        self.cursor = self.connection.cursor()

        # Create tables if they do not exist
        self.create_tables()

    def create_tables(self):
        """Create necessary tables in the database."""
        self.cursor.execute('''
            CREATE TABLE IF NOT EXISTS text_data (
                id INTEGER PRIMARY KEY AUTOINCREMENT,
                content TEXT
            )
        ''')
        
        self.cursor.execute('''
            CREATE TABLE IF NOT EXISTS links_data (
                id INTEGER PRIMARY KEY AUTOINCREMENT,
                text TEXT,
                url TEXT
            )
        ''')

        self.cursor.execute('''
            CREATE TABLE IF NOT EXISTS images_data (
                id INTEGER PRIMARY KEY AUTOINCREMENT,
                image BLOB,
                format TEXT
            )
        ''')

        self.cursor.execute('''
            CREATE TABLE IF NOT EXISTS tables_data (
                id INTEGER PRIMARY KEY AUTOINCREMENT,
                content TEXT  -- This stores table content as a string
            )
        ''')

        self.connection.commit()

    def save_text_data(self, text_data):
        """Save text data to SQL database."""
        for item in text_data:
            self.cursor.execute('INSERT INTO text_data (content) VALUES (?)', (item,))
        self.connection.commit()

    def save_links_data(self, links_data):
        """Save links data to SQL database."""
        for item in links_data:
            self.cursor.execute('INSERT INTO links_data (text, url) VALUES (?, ?)', (item['text'], item['url']))
        self.connection.commit()

    def save_images(self, images_data):
        """Save images data to SQL database."""
        for img_data in images_data:
            self.cursor.execute('INSERT INTO images_data (image, format) VALUES (?, ?)', 
                                (sqlite3.Binary(img_data['image_bytes']), img_data['image_format']))
        self.connection.commit()
        
    def save_tables(self, tables_data):
        """Save table data to SQL database as strings."""
        for i, table in enumerate(tables_data):
            try:
                # Convert table rows to strings, replacing None with empty strings
                table_str = '\n'.join([
                    ','.join([str(cell) if cell is not None else '' for cell in row]) for row in table
                ])
                # Save the table as text in the SQL database
                self.cursor.execute('INSERT INTO tables_data (content) VALUES (?)', (table_str,))
            except Exception as e:
                print(f"Error saving table {i + 1}: {e}")
        
        self.connection.commit()

    def close(self):
        """Close the database connection."""
        self.connection.close()


