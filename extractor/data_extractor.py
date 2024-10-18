import fitz  # PyMuPDF for PDF handling
import tabula  # For PDF table extraction
import docx  # For DOCX handling
from pptx import Presentation  # For PPTX handling
 
 
class DataExtractor:
    def __init__(self, loader):
        self.loader = loader
        self.file_path = loader.file_path
 
    def _load_file(self):
        """Helper method to load the appropriate file type."""
        if self.file_path.endswith('.pdf'):
            return fitz.open(self.file_path)
        elif self.file_path.endswith('.docx'):
            return self.loader.load_file()
        elif self.file_path.endswith('.pptx'):
            return self.loader.load_file()
        else:
            raise ValueError("Unsupported file format.")
 
    def _add_metadata(self, source, description, **kwargs):
        """Helper method to generate metadata."""
        metadata = {
            "source": source,
            "description": description
        }
        metadata.update(kwargs)
        return metadata
 
    def extract_text(self):
        """Extract text from the file."""
        text_data = []
        file = self._load_file()
 
        if self.file_path.endswith('.pdf'):
            for page_num in range(len(file)):
                page = file.load_page(page_num)
                text_data.append(page.get_text() or "")
            file.close()
 
        elif self.file_path.endswith('.docx'):
            text_data = [para.text for para in file.paragraphs]
 
        elif self.file_path.endswith('.pptx'):
            for slide in file.slides:
                text_data.extend([shape.text for shape in slide.shapes if hasattr(shape, "text")])
 
        return text_data
 
    def extract_images(self):
        """Extract images from the file, including metadata."""
        images_data = []
        file = self._load_file()
 
        if self.file_path.endswith('.pdf'):
            for page_num in range(len(file)):
                page = file.load_page(page_num)
                for img in page.get_images(full=True):
                    xref = img[0]
                    base_image = file.extract_image(xref)
                    images_data.append({
                        "image_bytes": base_image["image"],
                        "image_format": base_image["ext"],
                        "metadata": self._add_metadata(
                            self.file_path, f"Image from page {page_num + 1}",
                            page_number=page_num + 1, dimensions=(base_image["width"], base_image["height"])
                        )
                    })
            file.close()
 
        elif self.file_path.endswith('.docx'):
            for rel in file.part.rels:
                if "image" in file.part.rels[rel].target_ref:
                    images_data.append({
                        "image_bytes": file.part.rels[rel].target_part.blob,
                        "image_format": "png",  # Placeholder; determine format dynamically if possible
                        "metadata": self._add_metadata(self.file_path, "Image from DOCX file")
                    })
 
        elif self.file_path.endswith('.pptx'):
            for slide_num, slide in enumerate(file.slides):
                for shape in slide.shapes:
                    if hasattr(shape, "image"):
                        images_data.append({
                            "image_bytes": shape.image.blob,
                            "image_format": shape.image.ext,
                            "metadata": self._add_metadata(
                                self.file_path, f"Image from slide {slide_num + 1}",
                                slide_number=slide_num + 1
                            )
                        })
 
        return images_data
 
    def extract_links(self):
        """Extract hyperlinks from the file, including metadata."""
        links_data = []
        file = self._load_file()
 
        if self.file_path.endswith('.pdf'):
            for page_num in range(len(file)):
                page = file.load_page(page_num)
                for link in page.get_links():
                    if "uri" in link:
                        rect = link["from"]
                        links_data.append({
                            "text": link.get("title", link["uri"]),
                            "url": link["uri"],
                            "metadata": self._add_metadata(
                                self.file_path, f"Link from page {page_num + 1}",
                                page_number=page_num + 1, position=(rect.x0, rect.y0, rect.x1, rect.y1)
                            )
                        })
            file.close()
 
        elif self.file_path.endswith('.docx'):
            for rel in file.part.rels.values():
                if "hyperlink" in rel.target_ref:
                    links_data.append({
                        "text": rel.target_ref,
                        "url": rel.target_ref,
                        "metadata": self._add_metadata(self.file_path, "Hyperlink from DOCX file")
                    })
 
        elif self.file_path.endswith('.pptx'):
            for slide_num, slide in enumerate(file.slides):
                for shape in slide.shapes:
                    if hasattr(shape, "hyperlink") and shape.hyperlink.address:
                        links_data.append({
                            "text": getattr(shape, 'text', 'Hyperlinked shape'),
                            "url": shape.hyperlink.address,
                            "metadata": self._add_metadata(
                                self.file_path, f"Link from slide {slide_num + 1}",
                                slide_number=slide_num + 1
                            )
                        })
 
        return links_data
 
    def extract_tables(self):
        """Extract tables from the file, including metadata."""
        tables_data = []
        file = self._load_file()
 
        if self.file_path.endswith('.pdf'):
            try:
                tables = tabula.read_pdf(self.file_path, pages='all', multiple_tables=True)
                for table_index, table in enumerate(tables):
                    tables_data.append({
                        "content": table.values.tolist(),
                        "metadata": self._add_metadata(self.file_path, f"Table {table_index + 1} from PDF")
                    })
            except Exception as e:
                print(f"Error extracting tables from PDF: {e}")
 
        elif self.file_path.endswith('.docx'):
            for table_index, table in enumerate(file.tables):
                table_content = [[cell.text.strip() for cell in row.cells] for row in table.rows]
                tables_data.append({
                    "content": table_content,
                    "metadata": self._add_metadata(self.file_path, f"Table {table_index + 1} from DOCX")
                })
 
        elif self.file_path.endswith('.pptx'):
            for slide_num, slide in enumerate(file.slides):
                for shape in slide.shapes:
                    if hasattr(shape, "table"):
                        table_content = [[cell.text.strip() for cell in row.cells] for row in shape.table.rows]
                        tables_data.append({
                            "content": table_content,
                            "metadata": self._add_metadata(self.file_path, f"Table from slide {slide_num + 1}")
                        })
 
        return tables_data
 