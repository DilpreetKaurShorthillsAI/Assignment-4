import fitz  # PyMuPDF for PDF handling
import tabula  # For PDF table extraction
import docx  # For DOCX handling
from pptx import Presentation  # For PPTX handling

class DataExtractor:
    def __init__(self, loader):
        self.loader = loader
        self.file_path = loader.file_path

    def extract_text(self):
        """Extract text from the file."""
        text_data = []

        if self.file_path.endswith('.pdf'):
            # Extract text from PDF
            pdf_document = fitz.open(self.file_path)  # Open PDF file
            for page_num in range(len(pdf_document)):
                page = pdf_document.load_page(page_num)  # Load the page
                text_data.append(page.get_text() or "")  # Extract text from the page
            pdf_document.close()  # Close the document after reading
            return text_data

        elif self.file_path.endswith('.docx'):
            # Extract text from DOCX
            doc = self.loader.load_file()
            for paragraph in doc.paragraphs:
                text_data.append(paragraph.text + "\n")
            return text_data

        elif self.file_path.endswith('.pptx'):
            # Extract text from PPTX
            ppt = self.loader.load_file()
            for slide in ppt.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text_data.append(shape.text + "\n")
            return text_data

        else:
            raise ValueError("Unsupported file format for text extraction.")

    def extract_images(self):
        """Extract images from the file."""
        images_data = []

        if self.file_path.endswith('.pdf'):
            # PDF image extraction
            pdf_document = fitz.open(self.file_path)
            for page_num in range(len(pdf_document)):
                page = pdf_document.load_page(page_num)
                image_list = page.get_images(full=True)
                for img in image_list:
                    xref = img[0]
                    base_image = pdf_document.extract_image(xref)
                    image_bytes = base_image["image"]
                    image_ext = base_image["ext"]
                    width, height = base_image["width"], base_image["height"]
                    images_data.append({
                        "image_bytes": image_bytes,
                        "image_format": image_ext,
                        "page": page_num + 1,
                        "dimensions": (width, height)
                    })
            pdf_document.close()  # Close the document after extracting images

        elif self.file_path.endswith('.docx'):
            # DOCX image extraction
            doc = self.loader.load_file()
            for rel in doc.part.rels:
                if "image" in doc.part.rels[rel].target_ref:
                    image = doc.part.rels[rel].target_part.blob
                    images_data.append({
                        "image_bytes": image,
                        "image_format": "png"  # Default to png
                    })

        elif self.file_path.endswith('.pptx'):
            # PPTX image extraction
            ppt = self.loader.load_file()
            for slide in ppt.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "image"):
                        images_data.append({
                            "image_bytes": shape.image.blob,
                            "image_format": shape.image.ext
                        })

        return images_data

    def extract_links(self):
        """Extract hyperlinks from the file."""
        links_data = []

        if self.file_path.endswith('.pdf'):
            # PDF URL extraction
            pdf_document = fitz.open(self.file_path)
            for page_num in range(len(pdf_document)):
                page = pdf_document.load_page(page_num)
                links = page.get_links()
                for link in links:
                    if "uri" in link:
                        url = link["uri"]
                        text = link.get("title", url)  # Use title or fallback to URL
                        rect = link["from"]
                        links_data.append({
                            "text": text,  # Include text
                            "url": url,
                            "page": page_num + 1,
                            "position": {
                                "x0": rect.x0,
                                "y0": rect.y0,
                                "x1": rect.x1,
                                "y1": rect.y1
                            }
                        })
            pdf_document.close()  # Close the document after extracting links

        elif self.file_path.endswith('.docx'):
            # DOCX URL extraction
            doc = self.loader.load_file()
            for rel in doc.part.rels.values():
                if "hyperlink" in rel.target_ref:
                    links_data.append({"text": rel.target_ref, "url": rel.target_ref})

        elif self.file_path.endswith('.pptx'):
            # URLs are rare in PPTX, so skipping this part
            pass

        return links_data

    def extract_tables(self):
        """Extract tables from the file."""
        tables_data = []

        if self.file_path.endswith('.pdf'):
            # Extract tables from PDF using Tabula
            try:
                tables = tabula.read_pdf(self.file_path, pages='all', multiple_tables=True)
                for table in tables:
                    tables_data.append(table.values.tolist())  # Convert DataFrame to list of lists
            except Exception as e:
                print(f"Error extracting tables from PDF: {e}")
            return tables_data

        elif self.file_path.endswith('.docx'):
            # Extract tables from DOCX
            doc = self.loader.load_file()
            for table in doc.tables:
                table_content = [[cell.text for cell in row.cells] for row in table.rows]
                tables_data.append(table_content)
            return tables_data

        elif self.file_path.endswith('.pptx'):
            # Extract tables from PPTX
            ppt = self.loader.load_file()
            for slide in ppt.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "table"):
                        table = shape.table
                        table_content = [[cell.text for cell in row.cells] for row in table.rows]
                        tables_data.append(table_content)

        return tables_data
